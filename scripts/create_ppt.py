from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image
import os

prs = Presentation()
W = Inches(13.333)
H = Inches(7.5)
prs.slide_width = W
prs.slide_height = H

BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0x44, 0x44, 0x44)

# Project root (scripts/ -> parent)
BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TMP  = os.path.join(BASE, '_tmp_ppt')
os.makedirs(TMP, exist_ok=True)
OUT  = os.path.join(BASE, 'ppts', 'InfraredSR_Datasets_Overview.pptx')

FONT_NAME = 'Times New Roman'

def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

def add_title(slide, text, top=Inches(0.35)):
    tb = slide.shapes.add_textbox(Inches(0.6), top, Inches(12), Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(30); p.font.bold = True
    p.font.color.rgb = BLACK; p.font.name = FONT_NAME

def add_body(slide, text, left, top, width, height, fontSize=16, color=BLACK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(text.strip().split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(fontSize); p.font.color.rgb = color
        p.font.name = FONT_NAME; p.space_after = Pt(6)

def three_line_table(slide, data, left, top, col_widths, row_h=Inches(0.42), fontSize=14):
    rows, cols = len(data), len(data[0])
    tw = sum(col_widths)
    shape = slide.shapes.add_table(rows, cols, left, top, tw, row_h * rows)
    tbl = shape.table
    for i, w in enumerate(col_widths): tbl.columns[i].width = w
    for r in range(rows):
        tbl.rows[r].height = row_h
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = str(data[r][c])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(fontSize); para.font.name = FONT_NAME
                para.font.color.rgb = BLACK; para.font.bold = (r == 0)
                para.alignment = PP_ALIGN.CENTER
            tcPr = cell._tc.get_or_add_tcPr()
            for e in ('a:lnL','a:lnR','a:lnT','a:lnB','a:solidFill','a:noFill'):
                for ch in list(tcPr):
                    if ch.tag == qn(e): tcPr.remove(ch)

    thick  = {'w': 19000, 'color': '000000'}
    medium = {'w':  9500, 'color': '000000'}
    def sb(cell, edge, spec):
        tcPr = cell._tc.get_or_add_tcPr()
        ln = cell._tc.makeelement(qn(f'a:{edge}'), {'w': str(spec['w'])})
        sf = cell._tc.makeelement(qn('a:solidFill'), {})
        sc = cell._tc.makeelement(qn('a:srgbClr'), {'val': spec['color']})
        sf.append(sc); ln.append(sf); tcPr.append(ln)
    for c in range(cols):
        sb(tbl.cell(0, c), 'lnT', thick)
        sb(tbl.cell(0, c), 'lnB', medium)
        sb(tbl.cell(rows-1, c), 'lnB', thick)
    return tbl

def convert_img(src, dst, size=None):
    try:
        im = Image.open(src)
        if size: im = im.resize(size, Image.LANCZOS)
        if im.mode in ('RGBA','P'): im = im.convert('RGB')
        if im.mode != 'RGB': im = im.convert('RGB')
        im.save(dst); return dst
    except: return None

# ==================== SLIDE 1: 封面 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
add_title(s, "InfraredSR 项目数据集介绍", top=Inches(2.8))
add_body(s, "红外遥感方向五个数据集：超分辨率重建、航拍车辆检测、\n跨传感器图像翻译、跨光谱配准与可见光引导超分",
         Inches(0.6), Inches(4.2), Inches(12), Inches(1.0), fontSize=18, color=GRAY)

# ==================== SLIDE 2: 总览 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
add_title(s, "数据集总览")

d = [
    ["数据集", "任务", "模态", "分辨率 / 尺度", "规模"],
    ["Infrared SR", "红外图像超分辨率\n(×2, ×4)", "红外热成像\n(单通道灰度)", "HR 640×512\nLR×2 320×256 / LR×4 160×128", "7,121 对\n训练 4,983\n验证 1,069\n测试 1,069"],
    ["VEDAI", "航拍车辆检测\n(8 类)", "可见光 RGB + 红外\n(双模态配对)", "512×512\n1024×1024", "1,246 张标注图像\n3,706 个目标框\n5 折交叉验证"],
    ["OLI2MSI", "跨传感器图像翻译\n(Landsat-8 → Sentinel-2)", "多光谱卫星遥感\n(GeoTIFF)", "LR: 30m\nHR: 10m", "5,325 对\n训练 5,225\n测试 100"],
    ["CIDIS", "跨光谱图像配准\n(热红外 ↔ 可见光)", "热红外 + 可见光\n(配对)", "640×448\n(RGB)", "1,000 对\n训练 700\n验证 200\n测试 100"],
    ["CIDIS_guided\nx8 x16", "可见光引导红外\n超分辨率 (×8, ×16)", "热红外 + 可见光引导\n(配对)", "GT 640×448\nLR×8 80×56\nLR×16 40×28", "700 对训练\n200 对验证\n40 对测试"],
]
three_line_table(s, d, Inches(0.4), Inches(1.5),
    [Inches(1.7), Inches(2.3), Inches(2.3), Inches(2.8), Inches(2.3)],
    row_h=Inches(0.88), fontSize=13)

# ==================== SLIDE 3: Infrared SR ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
add_title(s, "Infrared SR — 红外图像超分辨率数据集")
info = (
    "• 任务：单帧红外图像超分辨率重建，包含 ×2 和 ×4 两种放大倍率\n"
    "• 格式：JPEG 灰度图像（8-bit 单通道），HR 与 LR 文件名一一对应\n"
    "• 目录：{trainsets, valsets, testsets} / {HR, LRx2, LRx4}\n"
    "• 分辨率：HR 640×512，LR×2 320×256，LR×4 160×128"
)
add_body(s, info, Inches(0.5), Inches(1.4), Inches(7.0), Inches(2.0), fontSize=16)

hr_p  = os.path.join(BASE, 'datasets/Data/trainsets/HR/0002.jpg')
lr2_p = os.path.join(BASE, 'datasets/Data/trainsets/LRx2/0002.jpg')
lr4_p = os.path.join(BASE, 'datasets/Data/trainsets/LRx4/0002.jpg')
y_img = Inches(3.6); x = Inches(0.8)
for label, src, w in [("HR 640×512", hr_p, Inches(3.0)), ("LR×2 320×256", lr2_p, Inches(1.5)), ("LR×4 160×128", lr4_p, Inches(0.75))]:
    dst = os.path.join(TMP, f'irsr_{label[:4]}.png')
    if convert_img(src, dst):
        s.shapes.add_picture(dst, x, y_img, width=w)
        add_body(s, label, x, y_img + w*0.55, w, Inches(0.3), fontSize=12, color=GRAY)
    x += w + Inches(0.6)

sd = [["集合", "HR", "LR×2", "LR×4"], ["训练", "4,983", "4,983", "4,983"], ["验证", "1,069", "1,069", "1,069"], ["测试", "1,069", "1,069", "1,069"]]
three_line_table(s, sd, Inches(8.5), Inches(1.4), [Inches(1.0)]*4, fontSize=13)
add_body(s, "示例图像 0002.jpg：同一场景在三种分辨率下的对照", Inches(0.5), Inches(6.7), Inches(8), Inches(0.4), fontSize=13, color=GRAY)

# ==================== SLIDE 4: VEDAI ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
add_title(s, "VEDAI — 航拍车辆检测数据集")
info = (
    "• 任务：航拍视角下的多类车辆目标检测，提供配对的可见光与红外图像\n"
    "• 格式：PNG 图像（*_co.png 为可见光，*_ir.png 为红外）+ YOLO 格式 TXT 标注\n"
    "• 版本：VEDAI (512×512) 和 VEDAI_1024 (1024×1024)，标注文件 1,246 个\n"
    "• 划分：5 折交叉验证（fold01~fold05.txt）\n"
    "• 标注：class_id  x_center  y_center  width  height（归一化坐标）"
)
add_body(s, info, Inches(0.5), Inches(1.4), Inches(12), Inches(2.2), fontSize=16)

co_p = os.path.join(BASE, 'datasets/dataset/VEDAI/images/00000000_co.png')
ir_p = os.path.join(BASE, 'datasets/dataset/VEDAI/images/00000000_ir.png')
y2 = Inches(3.8)
if convert_img(co_p, os.path.join(TMP,'vedai_co.png'), (320,320)):
    s.shapes.add_picture(os.path.join(TMP,'vedai_co.png'), Inches(1.0), y2, width=Inches(2.5))
    add_body(s, "可见光 RGB（00000000_co.png）", Inches(1.0), y2+Inches(2.6), Inches(2.5), Inches(0.3), fontSize=12, color=GRAY)
if convert_img(ir_p, os.path.join(TMP,'vedai_ir.png'), (320,320)):
    s.shapes.add_picture(os.path.join(TMP,'vedai_ir.png'), Inches(4.0), y2, width=Inches(2.5))
    add_body(s, "红外 IR（00000000_ir.png）", Inches(4.0), y2+Inches(2.6), Inches(2.5), Inches(0.3), fontSize=12, color=GRAY)

cd = [["类别", "轿车", "卡车", "皮卡", "拖拉机", "房车", "船", "摩托车", "其他"],
      ["数量", "955", "397", "307", "204", "190", "171", "105", "1,377"]]
three_line_table(s, cd, Inches(7.2), Inches(3.8), [Inches(0.9), Inches(0.65), Inches(0.65), Inches(0.65), Inches(0.7), Inches(0.65), Inches(0.55), Inches(0.8), Inches(0.7)], fontSize=12)

# ==================== SLIDE 5: OLI2MSI ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
add_title(s, "OLI2MSI — Landsat-8 → Sentinel-2 跨传感器图像翻译数据集")
info = (
    "• 任务：将 Landsat-8 OLI (30m) 图像翻译为 Sentinel-2 MSI (10m) 高质量图像\n"
    "    本质为带域迁移的超分辨率重建\n"
    "• 格式：GeoTIFF (.TIF) 多波段遥感图像\n"
    "• LR 输入：Landsat-8 OLI，30m 空间分辨率，单文件约 265 KB\n"
    "• HR 真值：Sentinel-2 MSI，10m 空间分辨率，单文件约 1.1 MB\n"
    "• 命名：L8_{场景ID}_{日期}_S2B_{日期}_T{瓦片号}_N{块号}.TIF\n"
    "• 5 个地理场景，2019-09-23 同日采集，按网格切分\n"
    "    场景 ID：126038, 126039, 126040, 126041, 126042"
)
add_body(s, info, Inches(0.5), Inches(1.4), Inches(8.2), Inches(4.8), fontSize=16)
od = [["集合", "图像对数量", "场景数", "说明"], ["训练集", "5,225", "5", "5 个场景按网格均匀分块"], ["测试集", "100", "5", "同一场景中取不重叠块"], ["合计", "5,325", "5", "Landsat-8 → Sentinel-2 一一配对"]]
three_line_table(s, od, Inches(9.0), Inches(1.4), [Inches(1.2), Inches(1.4), Inches(1.0), Inches(3.2)], fontSize=14)

# ==================== SLIDE 6: CIDIS ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
add_title(s, "CIDIS — 跨光谱图像配准数据集")
info = (
    "• 任务：热红外与可见光图像的跨光谱特征匹配与配准\n"
    "• 来源：Rivadeneira et al., 2024, \"Cross-Spectral Image Registration:\n"
    "    a Comparative Study and a New Benchmark Dataset\"\n"
    "• 格式：热红外 + 可见光配对图像，均为 640×448，RGB\n"
    "• 共 1,000 对：训练 700 / 验证 200 / 测试 100\n"
    "• 目录：{thermal, visible} / {train, val, test}\n"
    "• 场景涵盖多种真实环境，用于评估跨光谱配准算法性能"
)
add_body(s, info, Inches(0.5), Inches(1.4), Inches(7.5), Inches(3.5), fontSize=16)

# Sample images from CIDIS
cidis_th = os.path.join(BASE, 'datasets/CIDIS-dataset-main/dataset/thermal/train')
cidis_vi = os.path.join(BASE, 'datasets/CIDIS-dataset-main/dataset/visible/train')
th_files = sorted(os.listdir(cidis_th))
vi_files = sorted(os.listdir(cidis_vi))

y3 = Inches(4.0)
if th_files:
    if convert_img(os.path.join(cidis_th, th_files[0]), os.path.join(TMP,'cidis_th.png'), (280,196)):
        s.shapes.add_picture(os.path.join(TMP,'cidis_th.png'), Inches(8.5), Inches(1.4), width=Inches(2.0))
        add_body(s, "热红外", Inches(8.5), Inches(3.0), Inches(2.0), Inches(0.3), fontSize=12, color=GRAY)
if vi_files:
    if convert_img(os.path.join(cidis_vi, vi_files[0]), os.path.join(TMP,'cidis_vi.png'), (280,196)):
        s.shapes.add_picture(os.path.join(TMP,'cidis_vi.png'), Inches(11.0), Inches(1.4), width=Inches(2.0))
        add_body(s, "可见光", Inches(11.0), Inches(3.0), Inches(2.0), Inches(0.3), fontSize=12, color=GRAY)

cd_cidis = [["集合", "热红外", "可见光", "总计"],
            ["训练", "700", "700", "700 对"],
            ["验证", "200", "200", "200 对"],
            ["测试", "100", "100", "100 对"]]
three_line_table(s, cd_cidis, Inches(0.5), Inches(5.8), [Inches(1.2)]*4, fontSize=13)

# ==================== SLIDE 7: CIDIS_guided ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
add_title(s, "CIDIS_guided_x8x16 — 可见光引导红外超分辨率数据集")
info = (
    "• 任务：以可见光图像为引导，对热红外图像进行 ×8 和 ×16 超分辨率重建\n"
    "• 基于 CIDIS 数据集构建，添加合成退化（下采样）生成 LR 图像\n"
    "• 热红外 GT：640×448，LR×8：80×56，LR×16：40×28\n"
    "• 可见光引导图：640×448（与 GT 同尺寸，提供高频纹理引导）\n"
    "• 目录结构：{thermal, visible} / {train, val, test}\n"
    "    thermal 内含 {GT, LR_x8, LR_x16} 子目录\n"
    "• 训练 700 对 / 验证 200 对 / 测试 40 对（×8 和 ×16 各 40 张）"
)
add_body(s, info, Inches(0.5), Inches(1.4), Inches(7.5), Inches(4.0), fontSize=16)

# Sample: thermal GT, LRx8, LRx16, visible guide
th_gt_dir  = os.path.join(BASE, 'datasets/dataset_CIDIS_guided_x8x16/thermal/train/GT')
th_lr8_dir = os.path.join(BASE, 'datasets/dataset_CIDIS_guided_x8x16/thermal/train/LR_x8')
th_lr16_dir= os.path.join(BASE, 'datasets/dataset_CIDIS_guided_x8x16/thermal/train/LR_x16')
vis_dir    = os.path.join(BASE, 'datasets/dataset_CIDIS_guided_x8x16/visible/train')

gt_files  = sorted(os.listdir(th_gt_dir))
lr8_files = sorted(os.listdir(th_lr8_dir))
lr16_files= sorted(os.listdir(th_lr16_dir))
vis_files = sorted(os.listdir(vis_dir))

y4 = Inches(3.4)
items = [
    ("GT 640×448", th_gt_dir, gt_files, Inches(2.4), Inches(8.5), Inches(1.4)),
    ("可见光引导 640×448", vis_dir, vis_files, Inches(2.4), Inches(8.5), Inches(4.2)),
    ("LR×8 80×56", th_lr8_dir, lr8_files, Inches(0.7), Inches(8.5), Inches(5.1)),
    ("LR×16 40×28", th_lr16_dir, lr16_files, Inches(0.35), Inches(8.5), Inches(6.5)),
]
for label, d, files, w, lx, ty in items:
    if files:
        dst = os.path.join(TMP, f'cg_{label[:6]}.png')
        if convert_img(os.path.join(d, files[0]), dst):
            s.shapes.add_picture(dst, lx, ty, width=w)
            add_body(s, label, lx, ty+w*0.55, w, Inches(0.3), fontSize=11, color=GRAY)

cd_cg = [["集合", "热红外 (GT/LR×8/LR×16)", "可见光引导", "分辨率变化"],
         ["训练", "700×3 = 2,100", "700", "640×448 → 80×56 / 40×28"],
         ["验证", "200×3 = 600", "200", ""],
         ["测试", "40×2 = 80", "40", "×8 / ×16 各 40 张"]]
three_line_table(s, cd_cg, Inches(0.5), Inches(5.8), [Inches(1.2), Inches(2.8), Inches(1.8), Inches(2.5)], fontSize=12)

# ==================== SLIDE 8: 总结 ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
add_title(s, "总结与对比")

sm = [
    ["", "Infrared SR", "VEDAI", "OLI2MSI", "CIDIS", "CIDIS_guided"],
    ["任务", "单图超分辨率", "车辆目标检测", "跨传感器翻译", "跨光谱配准", "引导式超分辨率"],
    ["模态", "红外热成像\n(灰度)", "RGB + IR\n(双模态)", "L8 + S2\n多光谱", "热红外 +\n可见光", "热红外 +\n可见光引导"],
    ["分辨率", "HR 640×512\nLR 160×128", "512 / 1024", "LR 30m\nHR 10m", "640×448", "GT 640×448\nLR×16 40×28"],
    ["格式", "JPEG", "PNG + TXT", "GeoTIFF", "图像文件", "图像文件"],
    ["规模", "7,121 对", "1,246 张\n3,706 框", "5,325 对", "1,000 对", "700+200+40"],
    ["划分", "train/val/test", "5 折 CV", "train/test", "train/val/test", "train/val/test"],
]
three_line_table(s, sm, Inches(0.3), Inches(1.5),
    [Inches(1.2), Inches(2.0), Inches(2.2), Inches(2.0), Inches(2.0), Inches(2.4)],
    row_h=Inches(0.62), fontSize=12)

add_body(s, "五个数据集覆盖了红外遥感领域从像素级重建、目标检测、跨传感器增强、跨光谱配准到引导式超分的完整任务链。",
         Inches(0.5), Inches(6.5), Inches(12), Inches(0.4), fontSize=15, color=GRAY)

# ==================== Save ====================
prs.save(OUT)
print(f"Done: {OUT}")
